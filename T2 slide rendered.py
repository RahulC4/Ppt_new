import os
import uuid
import pythoncom
import win32com.client
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def export_slide_to_png(ppt_path, slide_index):
    """
    Uses PowerPoint COM to export a slide as PNG.
    """
    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True   # MUST be visible on enterprise laptops

    pres = powerpoint.Presentations.Open(ppt_path, WithWindow=True)

    slide = pres.Slides[slide_index]
    out_path = os.path.join(
        os.path.dirname(ppt_path),
        f"slide_{slide_index}_{uuid.uuid4().hex[:6]}.png"
    )

    slide.Export(out_path, "PNG", 1920, 1080)

    pres.Close()
    powerpoint.Quit()
    pythoncom.CoUninitialize()

    return out_path


def _is_editable_text_shape(shape):
    """
    Detect if the shape contains editable text.
    Ignore decorative text shapes inside groups.
    """
    if not shape.has_text_frame:
        return False

    text = shape.text.strip()
    if not text:
        return False

    # Ignore very short decorative labels (e.g., "1", "A", "→")
    if len(text) < 3 and shape.width < 100000:
        return False

    return True


def _extract_group_text_shapes(group_shape):
    """
    Auto-detect the main textboxes inside a group:

    - If only one textbox → return it
    - If multiple → return the one with longest text
    - Ignore decorative labels
    """
    text_shapes = []

    for shp in group_shape.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            text_shapes.extend(_extract_group_text_shapes(shp))
        elif _is_editable_text_shape(shp):
            text_shapes.append(shp)

    if not text_shapes:
        return []

    # Auto-detect the "main" editable shape → longest text wins
    main_shape = max(text_shapes, key=lambda s: len(s.text.strip()))

    return [main_shape]


def extract_slide_structure(ppt_path, slide_index):
    """
    Extract editable text shapes from the slide:
    - Titles
    - Body placeholders
    - Main text inside groups
    """
    prs = Presentation(ppt_path)
    slide = prs.slides[slide_index]

    editable_shapes = []
    idx = 0

    for shape in slide.shapes:

        # Case 1: simple textbox
        if _is_editable_text_shape(shape):
            shape_entry = {
                "shape_id": f"shape_{idx}",
                "text": shape.text.strip(),
                "placeholder": getattr(shape, "placeholder_format", None) is not None,
                "type": "title" if shape.is_placeholder and "title" in shape.name.lower()
                        else "body"
            }
            editable_shapes.append(shape_entry)
            idx += 1

        # Case 2: group shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shapes = _extract_group_text_shapes(shape)
            for shp in group_shapes:
                shape_entry = {
                    "shape_id": f"shape_{idx}",
                    "text": shp.text.strip(),
                    "placeholder": False,
                    "type": "body"
                }
                editable_shapes.append(shape_entry)
                idx += 1

    png_path = export_slide_to_png(ppt_path, slide_index)

    return {
        "slide_index": slide_index,
        "ppt_path": ppt_path,
        "png_path": png_path,
        "editable_shapes": editable_shapes
    }
