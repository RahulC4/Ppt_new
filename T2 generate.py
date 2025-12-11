# ==============================
# generate_ppt.py – Final Version (Option A)
# Text Replacement Per Shape + Bullet Preservation
# ==============================

import os
import uuid
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from utils import logger


def replace_text_in_shape(shape, new_text):
    """
    Replaces text inside a shape and preserves formatting + bullets.
    """
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    tf.clear()

    # Split on newlines → treat each as a bullet if original was bulleted
    lines = new_text.split("\n")

    for i, line in enumerate(lines):
        p = tf.add_paragraph()
        p.text = line.strip()
        p.font.size = Pt(18)

        # Preserve bullets if original had bullets
        try:
            if shape.text_frame.paragraphs[0].level > 0:
                p.level = shape.text_frame.paragraphs[0].level
            if shape.text_frame.paragraphs[0].bullet:
                p.font.bold = False
                p.level = 0
                p.bullet = True
        except:
            pass

        # Formatting
        p.alignment = PP_ALIGN.LEFT

    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def clone_slide(prs, source_slide):
    """
    Clone a slide by copying all its shapes into a new slide.
    This version keeps images, icons, formatting *because we don't rewrite relationships*.
    """
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    for shape in source_slide.shapes:
        el = shape.element
        new_el = el.clone()
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    return new_slide


def generate_presentation(selected_slides_data, user_answers):
    """
    Builds final PPT:
    - selected_slides_data: list of slide structures from slide_renderer
    - user_answers: dict { slide_index → { shape_id → "new text" } }
    """

    logger.info("Building final presentation from selected slides...")

    # Create new PPT
    new_prs = Presentation()

    for slide_struct in selected_slides_data:

        ppt_path = slide_struct["ppt_path"]
        slide_index = slide_struct["slide_index"]
        editable_shapes = slide_struct["editable_shapes"]

        # Load the original PPT
        prs = Presentation(ppt_path)
        source_slide = prs.slides[slide_index]

        # Clone the slide into new deck
        new_slide = clone_slide(new_prs, source_slide)

        # Get answers for this slide
        slide_answers = user_answers.get(str(slide_index), {})

        # Replace text per editable shape
        for shape_entry in editable_shapes:
            shape_id = shape_entry["shape_id"]
            original_text = shape_entry["text"]

            if shape_id not in slide_answers:
                continue

            new_text = slide_answers[shape_id]

            # Find matching shape in cloned slide
            for shp in new_slide.shapes:
                if shp.has_text_frame and shp.text.strip() == original_text:
                    replace_text_in_shape(shp, new_text)
                    break

    # Save new PPT
    out_path = os.path.join(
        "generated",
        f"ppt_{uuid.uuid4().hex[:6]}.pptx"
    )
    new_prs.save(out_path)

    logger.info(f"Generated PPT saved → {out_path}")
    return out_path
