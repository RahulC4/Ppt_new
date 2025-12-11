# generate_ppt.py
import os
import uuid
import tempfile
from copy import deepcopy
from pptx import Presentation
from utils import logger

def _clone_shape(shape, new_slide):
    """Deep copy PPT XML of shape to new slide."""
    try:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        return new_slide.shapes[-1]
    except Exception as e:
        logger.exception(f"Shape clone failed: {e}")
        return None

def clone_slide(src_ppt, slide_index, dst_ppt):
    """
    Clone slide from src_ppt[index] → dst_ppt
    Return cloned slide.
    """
    try:
        src = Presentation(src_ppt)
        src_slide = src.slides[slide_index]

        # Blank layout
        layout = dst_ppt.slide_layouts[6] if len(dst_ppt.slide_layouts) > 6 else dst_ppt.slide_layouts[0]
        new_slide = dst_ppt.slides.add_slide(layout)

        for shape in src_slide.shapes:
            _clone_shape(shape, new_slide)

        return new_slide

    except Exception as e:
        logger.exception(f"clone_slide failed: {e}")
        raise

def replace_text(slide, replacements: dict):
    """
    Replace text in slide using exact full-text matching.
    replacements = { original_text: new_text }
    """
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                original = shape.text.strip()
                if original in replacements and replacements[original] is not None:
                    new_text = replacements[original]
                    try:
                        tf = shape.text_frame
                        tf.clear()
                        tf.paragraphs[0].text = new_text
                    except:
                        shape.text = new_text
        except:
            continue

def generate_presentation_from_selected(selected_slides, answers_by_slide):
    """
    selected_slides: list of dicts → {ppt_path, slide_index, slide_id}
    answers_by_slide: { slide_id: { raw_replacements: {...} } }
    """

    out_ppt = Presentation()

    for slide_info in selected_slides:
        src_ppt = slide_info["ppt_path"]
        idx = slide_info["slide_index"]
        sid = slide_info["slide_id"]

        replacements = answers_by_slide.get(sid, {}).get("raw_replacements", {})

        new_slide = clone_slide(src_ppt, idx, out_ppt)

        replace_text(new_slide, replacements)

    out_path = os.path.join(tempfile.gettempdir(), f"generated_{uuid.uuid4().hex}.pptx")
    out_ppt.save(out_path)
    return out_path
