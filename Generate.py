# generate_ppt.py
import os
import tempfile
import uuid
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from azure_blob_utils import download_blob_to_file as unused_dl  # keep consistent name if present
from utils import logger

def _deepcopy_shape_to_slide(src_shape, dst_slide):
    """
    Deep-copy the shape XML element into the destination slide spTree.
    Returns the new shape object.
    """
    try:
        new_el = deepcopy(src_shape._element)
        # insert into the slide spTree
        spTree = dst_slide.shapes._spTree
        spTree.insert_element_before(new_el, 'p:extLst')
        # get the last added shape object
        return dst_slide.shapes[-1]
    except Exception as e:
        logger.exception(f"Failed to deepcopy shape: {e}")
        return None

def clone_slide_to_presentation(src_prs_path: str, src_index: int, dst_prs_obj: Presentation):
    """
    Clone a slide from source pptx (by index) into dst_prs_obj and return the new slide object.
    This performs a deep copy of shape xml to preserve design as much as possible.
    """
    try:
        src_prs = Presentation(src_prs_path)
        src_slide = src_prs.slides[src_index]

        # create a blank slide in destination (use layout 6 if available or first layout)
        try:
            blank_layout = dst_prs_obj.slide_layouts[6]
        except Exception:
            blank_layout = dst_prs_obj.slide_layouts[0]
        new_slide = dst_prs_obj.slides.add_slide(blank_layout)

        # copy slide-level background/theme by copying slide element (risky cross-file),
        # but copying shapes' xml is usually enough to preserve appearance
        for shape in src_slide.shapes:
            _deepcopy_shape_to_slide(shape, new_slide)

        return new_slide
    except Exception as e:
        logger.exception(f"Failed to clone slide index {src_index} from {src_prs_path}: {e}")
        raise

def _collect_text_map_from_slide(slide):
    """
    Returns a map { original_text_snippet_trimmed: [shape_index_in_slide, ...] }
    Used to identify shapes for replacement after cloning.
    """
    mapping = {}
    for i, shp in enumerate(slide.shapes):
        try:
            if hasattr(shp, "text") and shp.text and shp.text.strip():
                key = shp.text.strip()
                # use entire text as key (may be long)
                mapping.setdefault(key, []).append(i)
        except Exception:
            continue
    return mapping

def replace_text_in_slide(slide, replacements: dict):
    """
    Replace text in 'slide' based on the replacements dict.
    replacements mapping keys should be small substrings to match original text exactly or fully.
    Strategy:
      - Iterate text frames in the slide, if the full original text matches a key in replacements, replace.
      - If user asked 'keep original', that key will be omitted.
    """
    for shp in slide.shapes:
        try:
            if hasattr(shp, "text") and shp.text and shp.text.strip():
                orig = shp.text.strip()
                # If there's an exact replacement for this original text:
                if orig in replacements and replacements[orig] is not None:
                    new_text = replacements[orig]
                    # preserve paragraphs: clear then set new
                    try:
                        tf = shp.text_frame
                        tf.clear()
                        p = tf.paragraphs[0]
                        p.text = new_text
                    except Exception:
                        # fallback: set shape.text
                        shp.text = new_text
        except Exception:
            continue

def generate_presentation_from_selected(selected_slides_info, answers_by_slide, out_dir=None):
    """
    selected_slides_info: list of dicts with keys:
      - ppt_path (local path)
      - slide_index (int)
      - slide_id (string)
      - title/text/preview_image etc.

    answers_by_slide: dict keyed by slide_id; value is a dict:
      {
         "title": "New Title",
         "bullets": ["b1","b2",...],
         "raw_replacements": {"Original full text 1": "replacement text", ...}
      }

    Returns path to generated PPT.
    """
    if out_dir is None:
        out_dir = tempfile.gettempdir()

    out_prs = Presentation()

    for s_info in selected_slides_info:
        src_path = s_info["ppt_path"]
        idx = s_info["slide_index"]
        slide_id = s_info["slide_id"]

        # For making replacements mapping, prefer answers_by_slide[slide_id]['raw_replacements']
        repl = answers_by_slide.get(slide_id, {}).get("raw_replacements", {})

        # Clone slide into out_prs
        new_slide = clone_slide_to_presentation(src_path, idx, out_prs)

        # Replace text in the newly cloned slide
        replace_text_in_slide(new_slide, repl)

    # Save result
    fname = f"generated_{uuid.uuid4().hex[:8]}.pptx"
    out_path = os.path.join(out_dir, fname)
    out_prs.save(out_path)
    return out_path
